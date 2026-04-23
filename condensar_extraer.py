"""
condensar_extraer.py
Extrae texto de PDF / DOCX / PPTX / XLSX / MD — 100% local, 0 tokens LLM.

Uso:
    python condensar_extraer.py "<ruta_archivo>"

Salida: texto plano a stdout (UTF-8).
        Primera línea: "PAGINAS:<N>" para PDFs, "HOJAS:<N>" para Excel.
"""
import sys
import os
from pathlib import Path

sys.stdout.reconfigure(encoding='utf-8')
sys.stderr.reconfigure(encoding='utf-8')


def extraer_pdf(path: Path) -> str:
    try:
        import pdfplumber
    except ImportError:
        os.system(f'pip install pdfplumber -q')
        import pdfplumber

    paginas = []
    with pdfplumber.open(str(path)) as pdf:
        n_paginas = len(pdf.pages)
        print(f'PAGINAS:{n_paginas}', flush=True)
        for i, page in enumerate(pdf.pages):
            texto = page.extract_text()
            if texto:
                paginas.append(f'--- Página {i+1} ---\n{texto}')

    return '\n\n'.join(paginas)


def extraer_docx(path: Path) -> str:
    try:
        import docx
    except ImportError:
        os.system('pip install python-docx -q')
        import docx

    doc = docx.Document(str(path))
    parrafos = [p.text for p in doc.paragraphs if p.text.strip()]
    return '\n\n'.join(parrafos)


def extraer_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        os.system('pip install python-pptx -q')
        from pptx import Presentation

    prs = Presentation(str(path))
    diapositivas = []
    for i, slide in enumerate(prs.slides, 1):
        textos = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                textos.append(shape.text.strip())
        if textos:
            diapositivas.append(f'--- Diapositiva {i} ---\n' + '\n'.join(textos))

    return '\n\n'.join(diapositivas)


def extraer_excel(path: Path) -> str:
    """Extrae contenido de Excel (.xlsx/.xls) como tablas markdown."""
    try:
        import openpyxl
    except ImportError:
        os.system('pip install openpyxl -q')
        import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    hojas = []
    n_hojas = len(wb.sheetnames)
    print(f'HOJAS:{n_hojas}', flush=True)

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        filas = []
        for row in ws.iter_rows(values_only=True):
            # Convertir cada celda a string, None → vacío
            celdas = [str(c).strip() if c is not None else '' for c in row]
            # Saltar filas completamente vacías
            if not any(celdas):
                continue
            filas.append(celdas)

        if not filas:
            continue

        # Construir tabla markdown
        lineas = [f'--- Hoja: {sheet_name} ---']

        # Primera fila como header
        header = filas[0]
        lineas.append('| ' + ' | '.join(header) + ' |')
        lineas.append('| ' + ' | '.join(['---'] * len(header)) + ' |')

        # Resto como datos
        for fila in filas[1:]:
            # Ajustar cantidad de columnas al header
            while len(fila) < len(header):
                fila.append('')
            lineas.append('| ' + ' | '.join(fila[:len(header)]) + ' |')

        hojas.append('\n'.join(lineas))

    wb.close()
    return '\n\n'.join(hojas)


def extraer_md(path: Path) -> str:
    return path.read_text(encoding='utf-8', errors='replace')


def main():
    if len(sys.argv) < 2:
        print('ERROR: Uso: condensar_extraer.py "<ruta_archivo>"', file=sys.stderr)
        sys.exit(1)

    path = Path(sys.argv[1])
    if not path.exists():
        print(f'ERROR: Archivo no encontrado: {path}', file=sys.stderr)
        sys.exit(1)

    ext = path.suffix.lower()

    if ext == '.pdf':
        texto = extraer_pdf(path)
    elif ext == '.docx':
        texto = extraer_docx(path)
    elif ext == '.pptx':
        texto = extraer_pptx(path)
    elif ext in ('.xlsx', '.xls'):
        texto = extraer_excel(path)
    elif ext == '.md':
        texto = extraer_md(path)
    else:
        print(f'ERROR: Extensión no soportada: {ext}', file=sys.stderr)
        sys.exit(1)

    if not texto.strip():
        print(f'ERROR: No se pudo extraer texto de {path.name}', file=sys.stderr)
        sys.exit(2)

    print(texto)


if __name__ == '__main__':
    main()
